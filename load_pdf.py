#from langchain.document_loaders.pdf import  PyPDFDirectoryLoader
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.schema.document import Document
from dotenv import load_dotenv
import os
import logging
import time
import io
import fitz 
from office365.sharepoint.client_context import ClientContext
from office365.runtime.auth.user_credential import UserCredential
from office365.sharepoint.files.file import File
import urllib.parse
from pptx import Presentation
from openpyxl import load_workbook
from docx import Document as DocxDocument
from langchain_docling import DoclingLoader
logging.basicConfig(level=logging.INFO)
from docling.datamodel.base_models import DocumentStream
from docling.document_converter import DocumentConverter

import os


# Chargement des variables d'environnement
load_dotenv(dotenv_path='.config')

class LoadAndSplitDocuments:
    def __init__(self):
        self.sharepoint_url = "https://acipartnersparis2.sharepoint.com/sites/POC_RAG"
        self.username = "poc.rag@amitaconseil.com"
        self.password = "Juliepocrag2025"     

        # Connexion à SharePoint
        self.ctx = ClientContext(self.sharepoint_url).with_credentials(UserCredential(self.username, self.password))
        logging.info(f"Connexion à SharePoint avec {self.username}")

        # Vérification de la connexion
        try:
            web = self.ctx.web
            self.ctx.load(web)
            self.ctx.execute_query()
            logging.info(f"Connexion réussie à {web.properties['Title']}")
        except Exception as e:
            logging.error(f"Erreur de connexion : {e}")
            exit()

    # def get_files_from_sharepoint(self, extensions=[".pdf", ".pptx", ".xlsx", ".docx"], folder_url="/sites/POC_RAG/Documents%20partages"):
    #     """Récupère les fichiers depuis SharePoint (de manière récursive)"""
    #     all_files = []

    #     def explore_folder(folder):
    #         try:
    #             # Charger fichiers et sous-dossiers
    #             self.ctx.load(folder.files)
    #             self.ctx.load(folder.folders)
    #             self.ctx.execute_query()

    #             # Parcourir les fichiers
    #             for file in folder.files:
    #                 if any(file.properties["Name"].endswith(ext) for ext in extensions):
    #                     all_files.append(file)
    #                     logging.info(file.properties["Name"])

    #             # Appel récursif sur les sous-dossiers
    #             for subfolder in folder.folders:
    #                 explore_folder(subfolder)

    #         except Exception as e:
    #             logging.error(f"Erreur lors de l'exploration du dossier : {e}")

    #     try:
    #         root_folder = self.ctx.web.get_folder_by_server_relative_url(folder_url)
    #         explore_folder(root_folder)
    #     except Exception as e:
    #         logging.error(f"Erreur lors de l'accès au dossier racine : {e}")

    #     return all_files

    def get_files_from_sharepoint(self,extensions=[".pdf", ".pptx",".docx"],folder_url="/sites/POC_RAG/Documents%20partages"):
        """Récupère les fichiers depuis SharePoint."""
        all_files=[]
        def explore_folder(folder):
            files=folder.files
            subfolders=folder.folders
            self.ctx.load(files)
            self.ctx.load(subfolders)
            self.ctx.execute_query()

            for file in files:
                if any(file.properties["Name"].endswith(ext) for ext in extensions):
                    all_files.append(file)
                    logging.info(f"Found file: {file.properties['Name']}")

            for subfolder in subfolders:
                logging.info(f"Exploring subfolder: {subfolder.properties['Name']}")
                explore_folder(subfolder)  # appel récursif

        try:
            root_folder = self.ctx.web.get_folder_by_server_relative_url(folder_url)
            explore_folder(root_folder)

        except Exception as e:
            logging.error(f"Erreur lors de la récupération des fichiers SharePoint: {e}")
        return all_files
        # try:
        #     folder = self.ctx.web.get_folder_by_server_relative_url(folder_url)
        #     files = folder.files
        #     self.ctx.load(files)
        #     self.ctx.execute_query()
        #     for file in files :
        #         if any(file.properties["Name"].endswith(ext) for ext in extensions):
        #             all_files.append(file)
        #             logging.info(file.properties["Name"])
            
        #     for subfolder in folder.folders:
        #         print(1)
        #     # managed_files = [file for file in files if any(file.properties["Name"].endswith(ext) for ext in extensions)]
        #     # logging.info(f"{len(managed_files)} fichiers trouvés.")
        #     return all_files
        # except Exception as e:
        #     logging.error(f"Erreur lors de la récupération des fichiers : {e}")
        #     return []

    def load_files_from_memory(self, file):
        """Charge un fichier directement en mémoire depuis SharePoint et extrait son contenu."""
        try:
            # Obtient l'URL relative du fichier sur SharePoint
            server_relative_url = file.properties["ServerRelativeUrl"]

            # Charge le contenu du fichier depuis SharePoint en binaire
            file_content = File.open_binary(self.ctx, server_relative_url).content
            
            if not file_content:
                logging.error(f"Le fichier {file.properties['Name']} est vide ou introuvable.")
                return []
        
            file_name = file.properties["Name"].lower()
            file_extension = file_name.split('.')[-1]
            last_modified = file.properties.get("TimeLastModified", "Date inconnue").timestamp() #dernière modification
            logging.info(f"{last_modified}")
            if file_extension == "pdf":
                documents = []
                with io.BytesIO(file_content) as pdf_stream:
                    doc = fitz.open(stream=pdf_stream)  # Ouvre le PDF en mémoire
                    for page_num, page in enumerate(doc):
                        text = page.get_text()
                        documents.append(Document(
                                                page_content=text,
                                                metadata={
                                                    "numero_page": page_num + 1,  
                                                    "nom_fichier": file_name, 
                                                    "chemin_document": server_relative_url, 
                                                    "date_modification": last_modified,  
                                                    "dossier": os.path.dirname(server_relative_url)  
                                                    }))
                        #logging.info(f"PDF - Page {page_num + 1} Content: {text[:10]}...")
                return documents
            elif file_extension == "pptx":
                documents = []
                with io.BytesIO(file_content) as ppt_stream:
                    presentation = Presentation(ppt_stream)
                    for slide_num, slide in enumerate(presentation.slides):
                        slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                        documents.append(Document(
                                                page_content=slide_text,
                                                metadata={
                                                    "numero_page": slide_num + 1,  
                                                    "nom_fichier": file_name, 
                                                    "chemin_document": server_relative_url, 
                                                    "date_modification": last_modified,  
                                                    "dossier": os.path.dirname(server_relative_url)  
                                                }))
                        #logging.info(f"PPTX - Slide {slide_num + 1} Content: {slide_text[:10]}...")
                return documents
            # elif file_extension == "xlsx":
            #     documents = []
            #     with io.BytesIO(file_content) as excel_stream:
            #         workbook = load_workbook(excel_stream, data_only=True)
            #     for sheet in workbook.worksheets:
            #         sheet_text = ""
            #     for row in sheet.iter_rows():
            #         row_text = " | ".join(str(cell.value) if cell.value else "" for cell in row)
            #         sheet_text += row_text + "\n"
            #         documents.append(Document(
            #                                 page_content=sheet_text,
            #                                 metadata={
            #                                         "numero_page": sheet.title,  
            #                                         "nom_fichier": file_name, 
            #                                         "chemin_document": server_relative_url, 
            #                                         "date_modification": last_modified,  
            #                                         "dossier": os.path.dirname(server_relative_url)
            #                                 }))
            #         #logging.info(f"XLSX - Sheet {sheet.title} Content: {sheet_text[:10]}...")
            #     return documents
            elif file_extension == "docx":
                documents = []
                with io.BytesIO(file_content) as docx_stream:
                    doc = DocxDocument(docx_stream)
                    doc_text = "\n".join([para.text for para in doc.paragraphs])

                    documents.append(Document(
                                            page_content=doc_text,
                                            metadata={
                                                "numero_page": page_num + 1,  
                                                "nom_fichier": file_name, 
                                                "chemin_document": server_relative_url, 
                                                "date_modification": last_modified,  
                                                "dossier": os.path.dirname(server_relative_url)}
                                            ))
                    logging.info(f"DOCX - Content: {doc_text[:10]}...")
                return documents
            else:
                logging.warning(f"Format non pris en charge : {file_name}")
            return ""
        except Exception as e:
            logging.error(f"Erreur de lecture du fichier {file.properties['Name']} : {e}")
            return []
    


    def run_load_and_split_documents(self):
        """Charge et découpe les documents PDF sans les télécharger en local."""
        start_time = time.time()
        logging.info('Démarrage du traitement des documents')

        # Récupérer les fichiers depuis SharePoint
        files = self.get_files_from_sharepoint()
        all_chunks = []

        # Traiter chaque fichier
        for file in files:
            logging.info(f"Traitement du fichier : {file.properties['Name']}")
            documents = self.load_files_from_memory(file)  # Charger le contenu du fichier directement en mémoire
            chunks = self.split_documents(documents)
            all_chunks.extend(chunks)

        end_time = time.time()
        logging.info(f'Traitement terminé en {end_time - start_time} secondes')
        return all_chunks

    @staticmethod
    def split_documents(documents: list[Document]) -> list[Document]:
        """Découpe les documents en chunks."""
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,
            chunk_overlap=200,
            length_function=len,
            is_separator_regex=False,
        )
        return text_splitter.split_documents(documents)


#Exemple d'utilisation
if __name__ == "__main__":
    load_data = LoadAndSplitDocuments()
    
    # Récupère les fichiers depuis SharePoint
    files = load_data.get_files_from_sharepoint()

    # Charge et extrait le contenu du premier fichier
    if files:
        print("ok")
        # first_file = files[0]
        # documents = load_data.load_files_from_memory(first_file)
        # logging.info(f"Nombre de pages chargées : {len(documents)}")
    else:
        logging.warning("Aucun fichier trouvé.")


# if __name__ == '__main__':
#     load_and_split = LoadAndSplitDocuments()
#     chunks = load_and_split.run_load_and_split_documents()


