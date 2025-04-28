import logging
import os
import time

import streamlit as st
from dotenv import load_dotenv
from langchain.prompts import ChatPromptTemplate, PromptTemplate,MessagesPlaceholder, HumanMessagePromptTemplate
from langchain.retrievers.multi_query import MultiQueryRetriever
from langchain_community.vectorstores import FAISS
from langchain_chroma import Chroma
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_mistralai import ChatMistralAI,MistralAIEmbeddings
from langchain_core.messages import AIMessage, HumanMessage
from langchain.retrievers import MultiQueryRetriever 
from langchain.chains import ConversationalRetrievalChain,create_history_aware_retriever
from langchain.chains.combine_documents import create_stuff_documents_chain
#import hashlib
from datetime import datetime
from langchain.chains.query_constructor.base import AttributeInfo
from langchain.retrievers.self_query.base import SelfQueryRetriever
from langchain.retrievers import TimeWeightedVectorStoreRetriever
from langchain.retrievers.ensemble import EnsembleRetriever

from load_pdf import LoadAndSplitDocuments
from docx import Document as DocxDocument
from pptx import Presentation

from langchain.memory import ConversationBufferMemory

from langchain.docstore.document import Document
from langchain.document_loaders import PyPDFLoader, UnstructuredFileLoader
import tempfile
from bs4 import BeautifulSoup
from urllib.request import urlopen
import re
from langchain.chains import LLMChain
from pypdf.errors import EmptyFileError
import fitz
import io

load_dotenv(dotenv_path='.config')

def get_urls(message: str) -> list:
    url_pattern = r'https?://\S+|www\.\S+'
    match = re.findall(url_pattern, message) 
    if match :
        return match
    else : 
        return []
        
def get_url_content(url):
    page = urlopen(url)
    html = page.read().decode("utf-8")
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text().replace("\n", " ")

def format_docs(docs):
    """Fusionne une liste de Documents en texte structuré."""
    if not docs:
        return "*Aucun document.*"
    
    return "\n\n".join(
        f"[{doc.metadata.get('nom_fichier', 'document')}] (Page {doc.metadata.get('numero_page', 'N/A')})\n{doc.page_content}" 
        for doc in docs
    )

class InteractiveRAG:
    def __init__(self):
        #os.environ["OPENAI_API_KEY"] = os.getenv('OPENAI_API_KEY', '')
        os.environ["MISTRAL_API_KEY"] = os.getenv('MISTRAL_API_KEY', '')
        self.embedding_function = self.get_embedding_function()
        self._load_or_create_vector_db()
        # self.llm = ChatOpenAI(
        #     model="gpt-4o-mini",
        # #     model="gpt-3.5-turbo",
        #     temperature=0.7,
        #     max_tokens=None,
        # )
        self.llm = ChatMistralAI(
            model="mistral-small-latest",
            temperature=0.7,
            max_tokens=None,
            #random_seed=1
        )
        
        self.history_prompt = ChatPromptTemplate.from_messages([
            MessagesPlaceholder(variable_name="chat_history"),
            HumanMessagePromptTemplate.from_template("{input}")  
        ])
    
        self.query_prompt = PromptTemplate(
            input_variables=["question"],
            template="""
                Vous êtes un agent conversationnel pour les collaborateurs. 
                Reformulez cette question de deux manières différentes pour maximiser la récupération d'informations pertinentes.
                
                Question : {question}
            """
        )

        self.template_context_only = """Répondez à la question en vous appuyant uniquement sur le contexte suivant : {context}
            Question : {question}
            """
        
        self.template_mixed = """Répondez à la question en utilisant vos connaissances ainsi que le contexte suivant : {context}
            Question : {question}
            """
        
        self.prompt = ChatPromptTemplate.from_messages([
            ("system", 
            "Tu es un assistant expert.\n"
            "Réponds {mode_reponse}.\n\n"
            "Voici le contexte extrait de documents :\n{context}\n\n"
            "Documents uploadés :\n{uploaded_docs}\n\n"
            "Documents web :\n{web_doc}\n\n"
            "Base-toi sur cet historique de conversation également si besoin.\n"
            "Si une information est absente, incomplète ou incertaine :\n"
            "    • Ne tente pas de deviner.\n"
            "    • Indique clairement que tu ne disposes pas d'assez d'éléments.\n"
            "    • Propose, si possible, une reformulation ou une clarification à demander à l'utilisateur.\n"
            "**À la fin de ta réponse, liste explicitement les sources utilisées (noms de fichiers et le numéro de page).\n"
            "Si aucune source n'a été utilisée, indique clairement qu'aucune source n'a été utilisée dans ta réponse.**"
            ),
            MessagesPlaceholder(variable_name="chat_history_messages"),
            ("user", "{question}")
            ])

        self.metadata_field_info = [
                AttributeInfo(name="numero_page", description="Numéro de page", type="integer"),
                AttributeInfo(name="nom_fichier", description="Nom du fichier du document", type="string"),
                AttributeInfo(name="chemin_document", description="Chemin relatif serveur", type="string"),
                AttributeInfo(name="date_modification", description="Date de dernière modification (timestamp)", type="float"),
                AttributeInfo(name="dossier",description="Nom du dossier où est stocké le document",type="string")
                ]
        
        self.self_query_retriever = SelfQueryRetriever.from_llm(
            llm=self.llm,
            vectorstore=self.db,
            document_contents="Contenu du document",
            metadata_field_info=self.metadata_field_info,
            verbose=True
            )
        self.multi_query_retriever=MultiQueryRetriever.from_llm(
            retriever=self.db.as_retriever(search_kwargs={"k": 10}),
            llm=self.llm,
            prompt=self.query_prompt
        )
        self.time_weighted_retriever=TimeWeightedVectorStoreRetriever(
            vectorstore=self.db,
            other_score_keys=["date_modification"],
            k=10
        )
        self.ensemble_retriever = EnsembleRetriever(
            retrievers=[self.self_query_retriever, self.multi_query_retriever, self.time_weighted_retriever],
            weights=[0.4, 0.4, 0.2]
            )
        self.retriever = create_history_aware_retriever(
                        llm=self.llm,
                        retriever=self.ensemble_retriever,
                        prompt=self.history_prompt
                        )
        
    def _load_or_create_vector_db(self):
        #vector_db_path = "./faiss_index"
        #vector_db_path = "./chroma_index_openai"
        vector_db_path = "./chroma_index_mistral"

        if os.path.exists(vector_db_path):
            # Load existing vector store
            self.db = Chroma(persist_directory=vector_db_path, embedding_function=self.embedding_function)
            #self.db = FAISS.load_local(vector_db_path, self.embedding_function, allow_dangerous_deserialization=True)
        else:
            # Create and save vector store
            load_data = LoadAndSplitDocuments()
            document_chunks = load_data.run_load_and_split_documents()

            # Add sourcing metadata
            for doc in document_chunks:
                if "nom_fichier" not in doc.metadata:
                    doc.metadata["nom_fichier"] = "SharePoint"
                if "date_modification" not in doc.metadata:
                    doc.metadata["date_modification"] = 0 #datetime.now().isoformat().strftime("%Y-%m-%d %H:%M:%S")
                #doc.metadata["hash"] = hashlib.md5(doc.page_content.encode('utf-8')).hexdigest()
                
            # # Create vector store
            # self.db = FAISS.from_documents(document_chunks,
            #                                self.embedding_function)

            # # Persist vector store
            # self.db.save_local(vector_db_path)
     
            self.db = Chroma.from_documents(
            documents=document_chunks,
            embedding=self.embedding_function,
            persist_directory=vector_db_path
            )

    def update_vector_store_from_sharepoint(self): 

        directory = "./chroma_index_openai"
        #directory = "./chroma_index_mistral"
        load_data = LoadAndSplitDocuments()
        document_chunks = load_data.run_load_and_split_documents()

        # Ajout métadonnées manquantes
        for doc in document_chunks:
            if "nom_fichier" not in doc.metadata:
                doc.metadata["nom_fichier"] = "SharePoint"
            if "date_modification" not in doc.metadata:
                doc.metadata["date_modification"] = 0 #datetime.now().isoformat().strftime("%Y-%m-%d %H:%M:%S")
            #doc.metadata["hash"] = hashlib.md5(doc.page_content.encode('utf-8')).hexdigest()

        if os.path.exists(directory) and len(os.listdir(directory)) > 0:
            logging.info("Chargement du vector store existant")
            self.db = Chroma(persist_directory=directory, embedding_function=self.embedding_function)
            # Récupérer les documents existants avec leurs métadonnées
            existing_docs = self.db._collection.get(include=["metadatas", "documents"])
            existing_docs_metadata = {
                meta["nom_fichier"]: {"date_modification": meta["date_modification"], "hash": meta.get("hash", "")}
                for meta in existing_docs["metadatas"]
            }
            documents_to_add = []
            documents_to_update = []
            documents_to_remove = []

            #Documents à ajouter ou mettre à jour (date différente)
            for doc in document_chunks:
                source = doc.metadata["nom_fichier"]            
                hash_value = doc.metadata["hash"]
                last_modified = doc.metadata["date_modificationp"]
                if source not in existing_docs_metadata:
                    documents_to_add.append(doc)
                elif existing_docs_metadata[source]["date_modification"] != last_modified:
                    documents_to_update.append(doc)

            # Documents à supprimer
            existing_sources = {meta["nom_fichier"] for meta in existing_docs["metadatas"]}
            new_sources = {doc.metadata["nom_fichier"] for doc in document_chunks}
            documents_to_remove = list(existing_sources - new_sources)

            # Mettre à jour le vector store
            logging.info(f"Ajout : {len(documents_to_add)}, Mise à jour : {len(documents_to_update)}, Suppression : {len(documents_to_remove)}")

            if documents_to_add or documents_to_update or documents_to_remove:

                # Suppression des documents obsolètes
                if documents_to_remove:
                    self.db._collection.delete(where={"nom_fichier":{"$in":documents_to_remove}})
                    logging.info("Suppression")
                # Suppression ancienne version
                if documents_to_update:
                    self.db._collection.delete(where={"nom_fichier": {"$in": [doc.metadata["nom_fichier"] for doc in documents_to_update]}})

                # Ajout des nouveaux documents
                self.db.add_documents(documents_to_add + documents_to_update)
                logging.info("Ajout et mise à jour")
        else:
            logging.info("Création d'un nouveau vector store")
            # Créer un nouveau vector store
            self.db = Chroma.from_documents(
            documents=document_chunks,
            embedding=self.embedding_function,
            persist_directory=directory
            )
        #db.persist()
        logging.info(" Vector store enregistré localement avec succès !")

    def get_embedding_function(self):
        start_time = time.time()
        logging.info('get_embedding_function')
        #embeddings = OpenAIEmbeddings(model="text-embedding-ada-002")
        embeddings = MistralAIEmbeddings(model="mistral-embed")
        end_time = time.time()
        logging.info(f'get_embedding_function done in {end_time - start_time}')
        return embeddings
    
    def get_last_modified(self,file_path):
        timestamp = os.path.getmtime(file_path)
        return datetime.fromtimestamp(timestamp)

    def load_uploaded_docs(self,file_content, file_name, server_relative_url):
        with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
            tmp_file.write(file_content)
            tmp_file_path = tmp_file.name

        last_modified = self.get_last_modified(tmp_file_path)

        documents = []
        with io.BytesIO(file_content) as pdf_stream:
            doc = fitz.open(stream=pdf_stream)
            pdf_metadata = doc.metadata
        
            for page_num, page in enumerate(doc):
                text = page.get_text()
                documents.append(Document(
                                                page_content=text,
                                                metadata={
                                                    "numero_page": page_num + 1,  
                                                    "nom_fichier": file_name, 
                                                    "chemin_document": server_relative_url, 
                                                    "date_modification": last_modified,  
                                                    "dossier": os.path.dirname(server_relative_url)  
                                                    }))
    

        os.remove(tmp_file_path)
    
        return documents
    def load_uploaded_docs(self, file_content, file_name, server_relative_url):
        with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
            tmp_file.write(file_content)
            tmp_file_path = tmp_file.name
        
        last_modified = self.get_last_modified(tmp_file_path)

        file_extension = os.path.splitext(file_name)[1].lower().strip('.')
        
        documents = []

        if file_extension == "pdf":
            with io.BytesIO(file_content) as pdf_stream:
                doc = fitz.open(stream=pdf_stream)
                for page_num, page in enumerate(doc):
                    text = page.get_text()
                    documents.append(Document(
                        page_content=text,
                        metadata={
                            "numero_page": page_num + 1,
                            "nom_fichier": file_name,
                            "chemin_document": server_relative_url,
                            "date_modification": last_modified,
                            "dossier": os.path.dirname(server_relative_url)
                        }
                    ))

        elif file_extension == "pptx":
            with io.BytesIO(file_content) as ppt_stream:
                presentation = Presentation(ppt_stream)
                for slide_num, slide in enumerate(presentation.slides):
                    slide_text = "\n".join(
                        [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                    )
                    documents.append(Document(
                        page_content=slide_text,
                        metadata={
                            "numero_page": slide_num + 1,
                            "nom_fichier": file_name,
                            "chemin_document": server_relative_url,
                            "date_modification": last_modified,
                            "dossier": os.path.dirname(server_relative_url)
                        }
                    ))

        elif file_extension == "docx":
            with io.BytesIO(file_content) as docx_stream:
                doc = DocxDocument(docx_stream)
                doc_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
                documents.append(Document(
                    page_content=doc_text,
                    metadata={
                        "numero_page": 1,
                        "nom_fichier": file_name,
                        "chemin_document": server_relative_url,
                        "date_modification": last_modified,
                        "dossier": os.path.dirname(server_relative_url)
                    }
                ))

        else:
            raise ValueError(f"Unsupported file extension: {file_extension}")


        os.remove(tmp_file_path)
        
        return documents

    def run_rag_prompt(self, question: str,chat_history=None,uploaded_docs=None,use_model_knowledge=False,use_retrieved_context=True):

        start_time = time.time()
        logging.info('Run run_rag_prompt')
        logging.info('1. prompt')

        chat_history_messages = [
            HumanMessage(content=msg["content"]) if msg["role"] == "user" else AIMessage(content=msg["content"])
            for msg in chat_history[-3:]
        ] if chat_history else []

        #  Récupération des documents pertinents
        retrieved_docs=self.retriever.invoke({
            "input": question,
            "chat_history": chat_history_messages
        })
        sources = list(set([doc.metadata.get("nom_fichier", "Unknown") for doc in retrieved_docs]))

        uploaded_docs = uploaded_docs or []

        urls = get_urls(question)
        web_doc = []
        if urls:
            for url in urls:
                content = get_url_content(url)
                web_doc.append(Document(
                page_content=content,
                metadata={"nom_fichier": url}
                ))
        #  Création du prompt selon le mode
        #template = self.template_mixed if use_model_knowledge else self.template_context_only

        # template=self.template_context_with_history
        # prompt = ChatPromptTemplate.from_template(template)

        # question_answer_chain = create_stuff_documents_chain(self.llm, prompt=prompt)

        # result = question_answer_chain.invoke({
        # "chat_history_messages": chat_history_messages or [],
        # "context": (retrieved_docs or []) + (uploaded_docs or []) + (web_doc or []),
        # "uploaded_docs": uploaded_docs or [],
        # "web_doc": web_doc or [],
        # "question": question,
        # "use_model_knowledge": use_model_knowledge
        # })
        chain = LLMChain(
            llm=self.llm,
            prompt=self.prompt,
            verbose=True)
        
        result=chain.invoke( {"chat_history_messages": chat_history_messages or [],
                              "context": format_docs(retrieved_docs or []) if use_retrieved_context else "",
                              "uploaded_docs": format_docs(uploaded_docs or []),
                              "web_doc": format_docs(web_doc or []),
                              "question": question,
                              "mode_reponse": (
                                    "en s'appuyant uniquement sur le contexte fourni" 
                                    if not use_model_knowledge else 
                                    "en combinant ses propres connaissances avec le contexte fourni"
                                )
                            })


        # sources = list(set([doc.metadata.get("source", "Unknown") for doc in retrieved_docs]))[:2]

        # logging.info(f'1. prompt done {time.time() - start_time}')
        # logging.info('2. chain')
        # chain = (
        #     prompt
        #     | self.llm
        #     | StrOutputParser()
        # )

        # logging.info(f'2. chain done {time.time() - start_time}')
        # logging.info('3. result')
        # result = chain.invoke({
        #     "context": context_text,
        #     "question": question
        #})

#         sources = list(dict.fromkeys(
#           doc.metadata.get("source", "Unknown") for doc in retrieved_docs
#               ))
  
        logging.info(f'3. result done {time.time() - start_time}')
        logging.info(f'run_rag_prompt done {time.time() - start_time}')
        return {"response": result["text"], "sources": sources}


    def main(self):

        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []


        if "question" not in st.session_state:
            st.session_state.question = ""

        if "messages" not in st.session_state:
            st.session_state.messages = [
                {"role": "assistant", "content": "Bonjour ! Comment puis-je vous aider aujourd'hui ? 😊"}
        ]
        

        if "use_model_knowledge" not in st.session_state:
            st.session_state.use_model_knowledge = False
        if "use_retrieved_context" not in st.session_state:
            st.session_state.use_retrieved_context = True
    
        if not (st.session_state.use_model_knowledge or st.session_state.use_retrieved_context):
            st.session_state.use_retrieved_context = True
            st.rerun() 

        st.image("./image/img.png", width=200)
    
        with st.sidebar:

            st.header("Configuration")
            st.divider()

            #cols0 = st.columns(2)
            #with cols0[0]:
            st.toggle(
                    "Accès connaissances LLM", 
                    key="use_model_knowledge"
                )

            #with cols0[1]:
            st.toggle(
                    "Accès Sharepoint", 
                    key="use_retrieved_context"
                )
            
            if st.session_state.use_model_knowledge and st.session_state.use_retrieved_context:
                st.write("*Utilisation des connaissances et du contexte SharePoint*")
            elif st.session_state.use_model_knowledge:
                st.write("*Utilisation des connaissances uniquement*")
            elif st.session_state.use_retrieved_context:
                st.write("*Utilisation du contexte SharePoint uniquement*")

            if st.button("Nouvelle conversation"):
                    st.session_state.chat_history = []
                    st.session_state.messages.clear()
                    st.rerun() 
                #st.button("Clear Chat", on_click=lambda: st.session_state.messages.clear(), type="primary")

            st.header("Chargement de documents:")
            st.divider()
            # File upload input for RAG with documents
            uploaded_files = st.file_uploader(
                "📄 Charger un ou plusieurs document(s)", 
                type=["pdf", "pptx", "docx"], 
                accept_multiple_files=True,
                key="rag_docs",
                on_change=None 
                )

                    # Initialisation des states
            if "uploaded_files" not in st.session_state:
                st.session_state["uploaded_files"] = []
            if "uploaded_docs" not in st.session_state:
                st.session_state["uploaded_docs"] = []

            # Noms de fichiers uploadés actuellement visibles dans l'UI
            uploaded_file_names = [file.name for file in uploaded_files] if uploaded_files else []

            # Synchroniser : supprimer les fichiers retirés avec la croix
            files_to_delete = [
                file_name for file_name in st.session_state["uploaded_files"]
                if file_name not in uploaded_file_names
            ]

            for file_name in files_to_delete:
                st.session_state["uploaded_files"].remove(file_name)
                st.session_state["uploaded_docs"] = [
                    doc for doc in st.session_state["uploaded_docs"]
                    if doc.metadata.get("nom_fichier") != file_name
                ]

            # Ajouter les nouveaux fichiers
            if uploaded_files:
                new_files = [
                    file for file in uploaded_files
                    if file.name not in st.session_state["uploaded_files"]
                ]

                total_files_after_upload = len(st.session_state["uploaded_files"]) + len(new_files)
                if total_files_after_upload > 10:
                    st.warning("Vous avez atteint la limite de 10 documents.")
                else:
                    for file in new_files:
                        file.seek(0)
                        file_content = file.read()
                        file_name = file.name
                        server_relative_url = f"/uploaded_docs/{file_name}"

                        new_docs = self.load_uploaded_docs(file_content, file_name, server_relative_url)

                        # Ajout du fichier + ses pages
                        st.session_state["uploaded_files"].append(file_name)
                        st.session_state["uploaded_docs"].extend(new_docs)

            # # URL input for RAG with websites
            # st.text_input(
            #     "🌐 Insérer une URL", 
            #     placeholder="https://url.com",
            #     #on_change=load_url_to_db,
            #     key="rag_url",
            # )


        # Input field for user question
        #st.text_input("Welcome to GPT! Comment puis-je t'aider ? 😊")
        # Display previous chat messages
        for msg in st.session_state.messages:
            if msg["role"] == "user":
                with st.chat_message("user",avatar="👤"):
                    st.markdown(msg["content"])
            elif msg["role"] == "assistant":
                with st.chat_message("assistant",avatar="🤖"):
                    st.markdown(msg["content"])

        

        if prompt := st.chat_input("Your message"):
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user",avatar="👤"):
                st.markdown(prompt)

            # with st.chat_message("assistant",avatar="🤖"):
            #     message_placeholder = st.empty()
            #     full_response = ""

            messages = [HumanMessage(content=m["content"]) if m["role"] == "user" else AIMessage(content=m["content"]) for m in st.session_state.messages]

            with st.spinner("Récupération de la réponse..."):
                    # try:
                        response_data = self.run_rag_prompt(prompt, chat_history=st.session_state.messages,uploaded_docs=st.session_state.get("uploaded_docs", []), use_model_knowledge=st.session_state.use_model_knowledge,use_retrieved_context=st.session_state.use_retrieved_context)
                        full_response = response_data["response"]
                        with st.chat_message("assistant", avatar="🤖"):
                            st.markdown(full_response)
                        st.session_state.messages.append({"role": "assistant", "content": full_response})
                    # except Exception as e:
                    #     logging.error(f"Erreur lors de la génération de réponse : {e}")
                    #     full_response = "Désolé, une erreur s'est produite lors de la génération de la réponse."
                    #     message_placeholder.markdown(full_response)

                       

            # if not st.session_state.use_rag:
            #     st.write_stream(stream_llm_response(llm_stream, messages))
            # else:
            #     st.write_stream(stream_llm_rag_response(llm_stream, messages))
        # with st.chat_message("user"):
        #     st.markdown(prompt)
        # if question.strip():
        #     if st.session_state.query_submitted or question != "":
        #         st.session_state.chat_history.append({"role": "user", "message": question})
        #         # Display a progress bar
        #         with st.spinner('Génération de la réponse...'):
        #             progress_bar = st.progress(0)

        #             # Simulate response generation process
        #             for i in range(10):
        #                 time.sleep(0.1)  # Simulate time taken to generate response
        #                 progress_bar.progress((i + 1) * 10)

        #             # Generate answer using the RAG system
        #             result = self.run_rag_prompt(question=question,chat_history=st.session_state.chat_history,use_model_knowledge=use_model_knowledge)
        #             answer = result["response"]
        #             resources = result["resources"]

        #             # Add assistant's answer to chat history
        #             st.session_state.chat_history.append(
        #                 {"role": "assistant", "message": answer, "resources": resources})
        #         st.session_state.query_submitted = False

        # # Display chat history in reverse order (latest first)
        # st.write("### Conversation :")
        # pass
        # for chat in reversed(st.session_state.chat_history):
        #     if chat["role"] == "user":
        #         st.markdown(f"**Vous** : {chat['message']}")
        #     else:
        #         st.markdown(f"**GPT** : {chat['message']}")
        #         if "resources" in chat:
        #             st.markdown("**Ressources** :")
        #             for resource in chat['resources']:
        #                 st.markdown(f"- {resource}")


if __name__ == "__main__":
    rag = InteractiveRAG()
    rag.main()
